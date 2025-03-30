from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title & Application Introduction
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "[Your App Name]: Data Insights & Summaries at a Glance"
subtitle.text = "[Your Name/Company Name]\n\n[Your App Name] is a web application that transforms raw data into clear, actionable insights, summaries, and essential charts."

# Slide 2: Key Features & Analysis Capabilities
slide_layout = prs.slide_layouts[1]  # Title and Content Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Powerful Data Analysis & Visualization"
content.text = "- Automated Data Analysis and Summary Generation.\n" \
               "- Interactive Chart and Graph Creation (e.g., bar charts, line graphs, pie charts).\n" \
               "- Key Trend and Pattern Identification.\n" \
               "- Customizable Data Filtering and Segmenting.\n" \
               "- Exportable Reports and Visualizations."

# Slide 3: Impact & Benefits for Users
slide_layout = prs.slide_layouts[1]  # Title and Content Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Unlock Data-Driven Decisions"
content.text = "- Save Time and Effort in Data Analysis.\n" \
               "- Gain Clear and Concise Data Summaries.\n" \
               "- Identify Key Insights and Trends Quickly.\n" \
               "- Enhance Decision-Making with Visual Data.\n" \
               "- Increase Data Literacy."

# Save the presentation
pptx_filename = "/mnt/data/Data_Insights_Summary.pptx"
prs.save(pptx_filename)

pptx_filename
